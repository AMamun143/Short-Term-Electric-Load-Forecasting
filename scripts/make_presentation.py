#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(22)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = Path("results/Short_Term_Load_Forecasting_Presentation.pptx")
    metrics = root / "results" / "metrics.csv"

    prs = Presentation()
    add_title_slide(
        prs,
        "Short-Term Electric Load Forecasting",
        "Applied Soft Computing Project\nAbdullah Mamun",
    )

    slides = [
        ("Agenda", ["Problem", "Data", "Methods", "Results", "Takeaways"]),
        ("Problem", ["Need hourly load forecasts for grid operation", "High error means higher cost and risk"]),
        ("Goal", ["Build accurate and interpretable model", "Compare with strong baselines"]),
        ("Proposal Feedback Applied", ["Expanded literature view", "Added at least 3 baselines", "Kept ANFIS + PSO as main method"]),
        ("Dataset", ["AEP hourly load data", "Long time horizon and clear daily/weekly seasonality"]),
        ("Feature Engineering", ["Lag features: 1, 2, 3, 24, 48, 168 hours", "Calendar and cyclical time features"]),
        ("Train / Validation / Test Split", ["Time-based split: 70% / 15% / 15%", "No data leakage from future"]),
        ("Baseline 1: Persistence", ["Forecast = previous hour load", "Simple but strong benchmark"]),
        ("Baseline 2: Linear Regression", ["Captures linear dependence over lags", "Fast and interpretable"]),
        ("Baseline 3: MLP", ["Neural network baseline", "Can model non-linear behavior"]),
        ("Model 4: ANFIS-like", ["Neuro-fuzzy model with Gaussian memberships", "Combines fuzzy rules and data fitting"]),
        ("Model 5: ANFIS-PSO", ["PSO optimizes fuzzy premise parameters", "Consequent parameters solved by least squares"]),
        ("Evaluation Metrics", ["Point: MAE, RMSE, MAPE", "Interval: PICP and interval width"]),
        ("Reproducibility", ["One command: bash run.sh", "Downloads data, trains models, saves results and plots"]),
        ("Main Results", [f"Metrics file: {metrics}", "Best model selected by RMSE"]),
        ("Forecast Comparison Plot", ["Shows actual vs top 3 models", "Visual check for peak and valley tracking"]),
        ("Interval Forecast Plot", ["90% prediction interval from residual quantiles", "Shows uncertainty around point forecast"]),
        ("What Worked Well", ["PSO helped tune fuzzy premises", "Hybrid model improved stability"]),
        ("Limitations", ["Single-region data only", "No weather features in current version"]),
        ("Future Work", ["Add weather and holidays", "Try transfer learning across regions"]),
        ("Conclusion", ["Hybrid soft computing is promising for STLF", "Baselines are essential for fair comparison"]),
    ]

    # Keep exactly 20 slides total including title slide.
    for title, bullets in slides[:19]:
        add_bullet_slide(prs, title, bullets)

    prs.save(out)
    print(f"Saved presentation: {out.resolve()}")


if __name__ == "__main__":
    main()
